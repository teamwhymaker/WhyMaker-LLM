import os
import shutil
import openai
from dotenv import load_dotenv
from langchain_community.document_loaders import PyPDFLoader, UnstructuredImageLoader
import csv
from docx import Document as DocxDocument
from pptx import Presentation
try:
    from openpyxl import load_workbook
except Exception:
    load_workbook = None  # type: ignore
from langchain.schema import Document as TextDocument
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_chroma import Chroma
from typing import Optional, List
from langchain.chains import RetrievalQA, create_history_aware_retriever, create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.messages import HumanMessage, AIMessage
import json
from langchain_core.retrievers import BaseRetriever
from langchain_core.callbacks import CallbackManagerForRetrieverRun, BaseCallbackHandler
from langchain_core.documents import Document

load_dotenv()

# Set up OpenAI client
openai.api_key = os.getenv("OPENAI_API_KEY")
if not openai.api_key:
    # Allow container to start; requests will fail with clear error later
    print("Warning: OPENAI_API_KEY is not set.", flush=True)

# Helpers to read env at call time (avoids import-time capture in jobs)
def get_db_dir() -> str:
    return os.getenv("WHYMAKER_CHROMA_DIR", "/tmp/chroma_db")


def get_manifest_file() -> str:
    return os.getenv("WHYMAKER_MANIFEST_FILE", "/tmp/processed_files.json")
rag_chain = None


def _debug_log_dir(dir_path: str, label: str) -> None:
    try:
        import chromadb  # type: ignore
    except Exception:
        chromadb = None  # type: ignore
    if os.getenv("WHYMAKER_DEBUG_CHROMA", "false").lower() not in ("1", "true", "yes"):
        return
    try:
        entries = []
        if os.path.isdir(dir_path):
            for name in sorted(os.listdir(dir_path)):
                p = os.path.join(dir_path, name)
                kind = "dir" if os.path.isdir(p) else "file"
                entries.append(f"{name} ({kind})")
        print(
            f"[DEBUG] {label}: path={dir_path}, exists={os.path.exists(dir_path)}, chromadb_version={(getattr(chromadb, '__version__', 'unknown') if chromadb else 'unavailable')}, entries={entries}",
            flush=True,
        )
    except Exception as exc:
        print(f"[DEBUG] Failed to list {dir_path}: {exc}", flush=True)

def _ignore_manifest(dirpath: str, names: list[str]) -> list[str]:
    # Avoid copying manifest which can be written concurrently by the indexer
    return ["processed_files.json"] if "processed_files.json" in names else []


def _retrying_copytree(src: str, dst: str, attempts: int = 3, ignore_manifest: bool = True) -> None:
    """Copy src -> dst with retries to mitigate transient gcsfuse read errors.

    Retries on common gcsfuse messages like 'stale NFS file handle' or
    fuse ReadFileOp failures.
    """
    import time
    errors_to_retry = (
        "stale NFS file handle",
        "fuseops.ReadFileOp",
        "input/output error",
        "io: read/write on closed pipe",
    )
    last_exc: Exception | None = None
    for attempt in range(1, max(1, attempts) + 1):
        try:
            ignore = _ignore_manifest if ignore_manifest else None
            shutil.copytree(src, dst, dirs_exist_ok=True, ignore=ignore)
            return
        except Exception as exc:  # noqa: PERF203
            last_exc = exc
            msg = str(exc).lower()
            if any(token in msg for token in errors_to_retry) and attempt < attempts:
                sleep_s = min(2 ** attempt, 8)
                print(
                    f"Copy {src} -> {dst} failed (attempt {attempt}/{attempts}): {exc}. Retrying in {sleep_s}s...",
                    flush=True,
                )
                time.sleep(sleep_s)
                continue
            raise

def _open_vectorstore(embeddings: OpenAIEmbeddings) -> Chroma:
    """Open Chroma at DB_DIR; fall back to a writable /tmp store if unavailable.

    This prevents 500s when the mounted Chroma bucket is empty or read-only
    before the indexer job has populated it.
    """
    DB_DIR = get_db_dir()
    _debug_log_dir(DB_DIR, "Primary Chroma dir before open")
    try:
        return Chroma(persist_directory=DB_DIR, embedding_function=embeddings)
    except Exception as exc:
        print(f"Chroma open failed at {DB_DIR}: {exc}", flush=True)

        # If the primary path is a read-only GCSFuse mount, copy it to a local RW cache
        cache_dir = os.getenv("WHYMAKER_CHROMA_CACHE_DIR", "/tmp/chroma_db_cache")
        try:
            if os.path.isdir(DB_DIR):
                os.makedirs(cache_dir, exist_ok=True)
                db_manifest = os.path.join(DB_DIR, "processed_files.json")
                cache_manifest = os.path.join(cache_dir, "processed_files.json")

                should_copy = False
                # If cache missing or empty, copy
                if not os.listdir(cache_dir):
                    should_copy = True
                else:
                    # If DB manifest is newer than cache, refresh cache
                    try:
                        db_mtime = os.path.getmtime(db_manifest) if os.path.exists(db_manifest) else 0
                        cache_mtime = os.path.getmtime(cache_manifest) if os.path.exists(cache_manifest) else 0
                        if db_mtime > cache_mtime:
                            should_copy = True
                    except Exception:
                        should_copy = True

                if should_copy:
                    attempts = int(os.getenv("WHYMAKER_CACHE_COPY_RETRIES", "3"))
                    ignore_manifest = os.getenv("WHYMAKER_CACHE_IGNORE_MANIFEST", "true").lower() in ("1", "true", "yes")
                    _retrying_copytree(DB_DIR, cache_dir, attempts=attempts, ignore_manifest=ignore_manifest)

            # Try opening the cached copy
            _debug_log_dir(cache_dir, "Cached Chroma dir before open")
            return Chroma(persist_directory=cache_dir, embedding_function=embeddings)
        except Exception as cache_exc:
            print(f"Chroma cache open failed at {cache_dir}: {cache_exc}", flush=True)
            # Final fallback to an empty, ephemeral local store so the app stays up
            fallback_dir = os.getenv("WHYMAKER_FALLBACK_CHROMA_DIR", "/tmp/chroma_db_fallback")
            try:
                os.makedirs(fallback_dir, exist_ok=True)
            except Exception:
                pass
            _debug_log_dir(fallback_dir, "Fallback Chroma dir before open")
            return Chroma(persist_directory=fallback_dir, embedding_function=embeddings)

def _add_docs_with_token_safe_batches(vectordb: Chroma, docs: List[TextDocument], initial_batch_size: int = 64) -> None:
    """Add documents to Chroma in batches to avoid OpenAI max tokens per request.

    If an OpenAI BadRequestError with 'max_tokens_per_request' occurs, automatically
    reduce the batch size and retry. This ensures large corpora can be embedded
    without exceeding provider request limits.
    """
    if not docs:
        return
    batch_size = max(1, int(os.getenv("WHYMAKER_EMBED_BATCH", initial_batch_size)))
    i = 0
    while i < len(docs):
        end = min(i + batch_size, len(docs))
        slice_docs = docs[i:end]
        try:
            vectordb.add_documents(slice_docs)
            try:
                vectordb.persist()
            except Exception:
                pass
            i = end
        except Exception as e:
            msg = str(e)
            # Reduce batch size on token-limit errors and retry
            if "max_tokens_per_request" in msg or "Requested" in msg and "tokens" in msg:
                if batch_size == 1:
                    # Skip the problematic doc to avoid infinite loop
                    print(f"Skipping one oversized chunk due to token limits at index {i}", flush=True)
                    i += 1
                else:
                    batch_size = max(1, batch_size // 2)
                    print(f"Token limit hit. Reducing embed batch size to {batch_size} and retrying...", flush=True)
            else:
                # Log and skip this slice to keep job moving
                print(f"Embedding error for batch {i}:{end} -> {e}", flush=True)
                i = end

def process_documents(folder_path="uploads"):
    """
    Incrementally processes documents from folder_path, chunking and adding
    them to the Chroma vector store in a memory-safe way. Persists after
    each file and updates MANIFEST_FILE to avoid re-processing.
    """
    print("Processing documents...", flush=True)

    # Load manifest of already processed files (filename -> mtime)
    MANIFEST_FILE = get_manifest_file()
    if os.path.exists(MANIFEST_FILE):
        with open(MANIFEST_FILE, "r") as f:
            processed = json.load(f)
    else:
        processed = {}

    supported_exts = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".pptx": "pptx",
        ".jpg": "image",
        ".jpeg": "image",
        ".png": "image",
        ".svg": "image",
        ".txt": "text",
        ".md": "text",
        ".csv": "csv",
        ".xlsx": "xlsx",
    }

    # Prepare vector store once
    embeddings = OpenAIEmbeddings(model="text-embedding-ada-002")
    vectordb = _open_vectorstore(embeddings)
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)

    new_files = 0
    for current_dir, _, files in os.walk(folder_path):
        # Sort for deterministic ordering
        for filename in sorted(files):
            if filename.startswith((".", "~$")):
                continue
            ext = os.path.splitext(filename.lower())[1]
            file_path = os.path.join(current_dir, filename)
            rel_key = os.path.relpath(file_path, folder_path)
            if ext not in supported_exts:
                print(f"  - Skipping unsupported: {rel_key} ({ext})", flush=True)
                continue
            try:
                mtime = os.path.getmtime(file_path)
            except FileNotFoundError:
                continue
            if processed.get(rel_key) == mtime:
                continue

            print(f"  - Ingesting {rel_key}", flush=True)

            # Load per-file docs
            try:
                if supported_exts[ext] == "pdf":
                    loader = PyPDFLoader(file_path)
                    docs = loader.load()
                elif supported_exts[ext] == "docx":
                    docs = []
                    doc = DocxDocument(file_path)
                    for para in doc.paragraphs:
                        text = para.text.strip()
                        if text:
                            docs.append(TextDocument(page_content=text, metadata={"source": file_path}))
                    for table in getattr(doc, "tables", []):
                        table_text = "\n".join([" | ".join([cell.text.strip() for cell in row.cells]) for row in table.rows])
                        if table_text:
                            docs.append(TextDocument(page_content=table_text, metadata={"source": file_path, "type": "table"}))
                elif supported_exts[ext] == "pptx":
                    docs = []
                    pres = Presentation(file_path)
                    for i, slide in enumerate(pres.slides):
                        for shape in slide.shapes:
                            if getattr(shape, "has_table", False):
                                table = shape.table  # type: ignore
                                table_text = "\n".join([" | ".join([cell.text.strip() for cell in row.cells]) for row in table.rows])
                                if table_text:
                                    docs.append(TextDocument(page_content=table_text, metadata={"source": file_path, "slide": i, "type": "table"}))
                            elif getattr(shape, "has_text_frame", False):
                                text = shape.text  # type: ignore
                                if text:
                                    docs.append(TextDocument(page_content=text, metadata={"source": file_path, "slide": i}))
                elif supported_exts[ext] == "image":
                    loader = UnstructuredImageLoader(file_path)
                    docs = loader.load()
                elif supported_exts[ext] == "text":
                    try:
                        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                            text = f.read()
                    except Exception:
                        text = ""
                    docs = [TextDocument(page_content=text, metadata={"source": file_path})]
                elif supported_exts[ext] == "csv":
                    table_lines = []
                    try:
                        with open(file_path, "r", encoding="utf-8", errors="ignore", newline="") as f:
                            sample = f.read(4096)
                            f.seek(0)
                            try:
                                dialect = csv.Sniffer().sniff(sample)
                            except Exception:
                                dialect = csv.excel
                            reader = csv.reader(f, dialect)
                            for row in reader:
                                table_lines.append(" | ".join([str(cell) for cell in row]))
                    except Exception:
                        pass
                    docs = [TextDocument(page_content="\n".join(table_lines), metadata={"source": file_path, "type": "table"})]
                elif supported_exts[ext] == "xlsx":
                    docs = []
                    if load_workbook is None:
                        print("openpyxl not installed; skipping xlsx: " + file_path, flush=True)
                    else:
                        try:
                            wb = load_workbook(file_path, read_only=True, data_only=True)
                            for sheet in wb.worksheets:
                                rows_text = []
                                for row in sheet.iter_rows(values_only=True):
                                    cells = ["" if v is None else str(v) for v in row]
                                    rows_text.append(" | ".join(cells))
                                if rows_text:
                                    docs.append(TextDocument(
                                        page_content="\n".join(rows_text),
                                        metadata={"source": file_path, "sheet": sheet.title, "type": "table"},
                                    ))
                        except Exception as e:
                            print(f"Failed to read xlsx {file_path}: {e}", flush=True)
                else:
                    docs = []
            except Exception as load_exc:
                print(f"  - Error loading {rel_key}: {load_exc}", flush=True)
                continue

            # Split and add with token-safe batching
            chunks = text_splitter.split_documents(docs)
            if chunks:
                _add_docs_with_token_safe_batches(vectordb, chunks)
            processed[rel_key] = mtime
            new_files += 1

        # Save manifest incrementally
        try:
            os.makedirs(os.path.dirname(MANIFEST_FILE), exist_ok=True)
        except Exception:
            pass
        with open(MANIFEST_FILE, "w") as f:
            json.dump(processed, f)

    if new_files == 0:
        print("No new documents to process.")
    else:
        print(f"  - Ingested {new_files} files.", flush=True)

class CompositeRetriever(BaseRetriever):
    """
    A custom retriever that combines results from a base retriever
    with a list of extra, in-memory documents to make it LCEL-compatible.
    """
    base_retriever: BaseRetriever
    extra_docs: List[Document]

    def _get_relevant_documents(
        self, query: str, *, run_manager: CallbackManagerForRetrieverRun
    ) -> List[Document]:
        """
        Get documents from the base retriever and combine with extra docs.
        """
        base_docs = self.base_retriever.get_relevant_documents(query, callbacks=run_manager.get_child())
        return base_docs + self.extra_docs


def setup_rag_chain(
    model_name: str = "o4-mini",
    retriever: BaseRetriever = None,
    streaming: bool = False,
    llm_callbacks: Optional[List[BaseCallbackHandler]] = None,
):
    """
    Sets up a conversational LangChain RAG chain for a given model.
    """
    print(f"Setting up RAG chain with model: {model_name}...", flush=True)
    if not retriever:
      embeddings = OpenAIEmbeddings(model="text-embedding-ada-002")
      vectordb = Chroma(persist_directory=get_db_dir(), embedding_function=embeddings)
      retriever = vectordb.as_retriever()
    
    llm = ChatOpenAI(model=model_name, streaming=streaming, callbacks=llm_callbacks)

    # Contextualize question prompt
    contextualize_q_system_prompt = (
        "Given a chat history and the latest user question "
        "which might reference context in the chat history, "
        "formulate a standalone question which can be understood "
        "without the chat history. Do NOT answer the question, "
        "just reformulate it if needed and otherwise return it as is."
    )
    contextualize_q_prompt = ChatPromptTemplate.from_messages(
        [
            ("system", contextualize_q_system_prompt),
            MessagesPlaceholder("chat_history"),
            ("human", "{input}"),
        ]
    )
    history_aware_retriever = create_history_aware_retriever(
        llm, retriever, contextualize_q_prompt
    )

    # Answering prompt
    qa_system_prompt = (
        "You are a world-class business and educational assistant, specifically tailored for the WhyMaker team. "
        "Your primary goal is to help WhyMaker staff create high-quality materials, including sales scripts, "
        "marketing collateral, lesson plans, and more.\\n\\n"
        "To answer the user's request, synthesize information from two sources: "
        "1. The provided internal WhyMaker documents (retrieved context below). "
        "2. Your general knowledge for broader context and information not available in the documents.\\n\\n"
        "CRITICAL INSTRUCTIONS:\\n"
        "- Provide comprehensive, well-structured, and clear responses. Do not be overly brief.\\n"
        "- Use Markdown formatting (like ### Headers, * Bullet Points, and **bold text**) to improve readability "
        "and ensure your answers are easy to interpret.\\n"
        "- If the provided context doesn't contain a specific answer, clearly state that the information isn't "
        "in WhyMaker's documents. Then, provide the best possible answer based on your general knowledge, "
        "while noting it may not be specific to WhyMaker.\\n\\n"
        "CONTEXT:\\n{context}"
    )
    qa_prompt = ChatPromptTemplate.from_messages(
        [
            ("system", qa_system_prompt),
            MessagesPlaceholder("chat_history"),
            ("human", "{input}"),
        ]
    )

    question_answer_chain = create_stuff_documents_chain(llm, qa_prompt)

    # Return the chain instead of setting a global variable
    return create_retrieval_chain(history_aware_retriever, question_answer_chain)


def query_rag(
    question: str,
    chat_history: Optional[list] = None,
    model_name: str = "o4-mini",
    extra_docs: Optional[list] = None,
):
    """
    Queries the RAG system with a question, history, and model name.
    """
    embeddings = OpenAIEmbeddings(model="text-embedding-ada-002")
    vectordb = _open_vectorstore(embeddings)
    base_retriever = vectordb.as_retriever()

    if extra_docs:
        retriever = CompositeRetriever(base_retriever=base_retriever, extra_docs=extra_docs)
    else:
        retriever = base_retriever

    # Create a new chain for each query, configured with the selected model and retriever
    rag_chain = setup_rag_chain(model_name, retriever=retriever, streaming=False)

    if chat_history is None:
        chat_history = []
        
    result = rag_chain.invoke({"input": question, "chat_history": chat_history})
    
    answer = result.get("answer", "I don't have enough information to answer.")

    # --- Debugging: Print retrieved source documents ---
    # This part of the original code was not included in the new_code,
    # so it is removed to match the new_code's structure.
    # if result.get("source_documents"):
    #     print("\nRetrieved documents:", flush=True)
    #     for doc in result.get("source_documents"):
    #         source = doc.metadata.get("source")
    #         content_preview = doc.page_content[:100] + "..." if len(doc.page_content) > 100 else doc.page_content
    #         print(f"  - From: {source}\n    Preview: \"{content_preview}\"", flush=True)
    # print("-------------------------\n", flush=True)
    
    return answer, None # No title generation for now


def main():
    """
    Main function to run the RAG CLI.
    """
    print("Starting the RAG system...", flush=True)
    
    # Ingest documents (incremental): process new or updated files
    process_documents()

    # setup_rag_chain() # This line is removed as per the edit hint
    
    print("\nRAG system is ready. Ask your questions!", flush=True)
    print("Type 'exit' to quit.", flush=True)
    
    while True:
        user_question = input("\n> ")
        if user_question.lower() == 'exit':
            break
        
        answer = query_rag(user_question)
        print(f"\nAssistant: {answer}", flush=True)

if __name__ == "__main__":
    main() 